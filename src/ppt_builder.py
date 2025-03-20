from pptx import Presentation
from pptx.util import Inches
import yaml
import os

class PPTBuilder:
    def __init__(self):
        with open('config/paths.yaml') as f:
            self.path_config = yaml.safe_load(f)
    
    def build_from_mindmap(self, json_data):
        """根据结构化数据生成PPT"""
        prs = Presentation(self.path_config['default_paths']['ppt_template'])
        
        # 添加标题页
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = json_data['title']
        subtitle.text = "自动生成报告"
        
        # 添加内容页
        for section in json_data['sections']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = section['title']
            content.text = '\n'.join(section['points'])
            
        prs.save(self.path_config['default_paths']['ppt_output'])
